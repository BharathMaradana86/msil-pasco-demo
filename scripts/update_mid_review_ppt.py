from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Sequence

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import PP_PLACEHOLDER
from pptx.enum.text import PP_ALIGN
from pptx.util import Pt


@dataclass(frozen=True)
class SlideSpec:
    title: str
    bullets: Sequence[str]
    subtitle_lines: Optional[Sequence[str]] = None  # only used for title slide


FONT_NAME = "Aptos"
TITLE_SIZE = 44
BODY_SIZE = 24
SUBTITLE_SIZE = 20

# Theme (simple, modern, high-contrast)
ACCENT = RGBColor(0, 170, 190)        # teal
INK = RGBColor(17, 24, 39)           # near-black
MUTED = RGBColor(71, 85, 105)        # slate
BG = RGBColor(248, 250, 252)         # very light gray
CARD_BG = RGBColor(255, 255, 255)    # white
BORDER = RGBColor(226, 232, 240)     # light border

MARGIN_X = Pt(44)
HEADER_H = Pt(72)
FOOTER_H = Pt(28)
GUTTER = Pt(18)


def _slide_w(prs: Presentation) -> int:
    return int(prs.slide_width)


def _slide_h(prs: Presentation) -> int:
    return int(prs.slide_height)


def _px(value):
    # pptx uses EMU under the hood; Pt already returns EMU.
    return value


def _apply_theme_background(slide, prs: Presentation) -> None:
    # Soft background fill
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _clear_all_textboxes(slide) -> None:
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            shape.text_frame.clear()


def _add_header(slide, prs: Presentation, title: str, *, show_tag: bool = True) -> None:
    # Header band
    header = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE (avoid extra import)
        0,
        0,
        prs.slide_width,
        HEADER_H,
    )
    header.fill.solid()
    header.fill.fore_color.rgb = CARD_BG
    header.line.color.rgb = BORDER
    header.line.width = Pt(1)

    # Accent strip
    strip = slide.shapes.add_shape(1, 0, 0, Pt(10), HEADER_H)
    strip.fill.solid()
    strip.fill.fore_color.rgb = ACCENT
    strip.line.fill.background()

    # Title text
    tx = slide.shapes.add_textbox(MARGIN_X, Pt(18), prs.slide_width - (MARGIN_X * 2), Pt(40))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.name = FONT_NAME
    p.font.size = Pt(30)
    p.font.bold = True
    p.font.color.rgb = INK

    if show_tag:
        tag = slide.shapes.add_textbox(prs.slide_width - Pt(210), Pt(22), Pt(170), Pt(26))
        tf2 = tag.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        p2.text = "Mid Project Review"
        p2.font.name = FONT_NAME
        p2.font.size = Pt(12)
        p2.font.bold = True
        p2.font.color.rgb = MUTED
        p2.alignment = PP_ALIGN.RIGHT


def _add_footer(slide, prs: Presentation, text: str) -> None:
    footer = slide.shapes.add_shape(
        1,
        0,
        prs.slide_height - FOOTER_H,
        prs.slide_width,
        FOOTER_H,
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = CARD_BG
    footer.line.color.rgb = BORDER
    footer.line.width = Pt(1)

    tx = slide.shapes.add_textbox(MARGIN_X, prs.slide_height - FOOTER_H + Pt(6), prs.slide_width - (MARGIN_X * 2), Pt(18))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(10)
    p.font.color.rgb = MUTED


def _add_card(slide, left, top, width, height):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = CARD_BG
    card.line.color.rgb = BORDER
    card.line.width = Pt(1)
    return card


def _add_card_title(slide, left, top, width, text: str):
    tx = slide.shapes.add_textbox(left, top, width, Pt(24))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = text
    p.font.name = FONT_NAME
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = INK
    return tx


def _add_bullets_in_box(slide, left, top, width, height, bullets: Sequence[str], *, bullet_symbol: str = "•") -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.margin_left = Pt(12)
    tf.margin_right = Pt(10)
    tf.margin_top = Pt(6)
    tf.margin_bottom = Pt(6)

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"{bullet_symbol} {b}"
        p.level = 0
        p.font.name = FONT_NAME
        p.font.size = Pt(16)
        p.font.color.rgb = MUTED
        p.space_after = Pt(6)


def _render_title_slide(slide, prs: Presentation, spec: SlideSpec) -> None:
    _apply_theme_background(slide, prs)
    _clear_all_textboxes(slide)

    # Center card
    card_w = prs.slide_width - (MARGIN_X * 2)
    card_h = Pt(320)
    card_left = MARGIN_X
    card_top = Pt(140)
    _add_card(slide, card_left, card_top, card_w, card_h)

    # Accent pill
    pill = slide.shapes.add_shape(1, card_left + Pt(18), card_top + Pt(18), Pt(140), Pt(26))
    pill.fill.solid()
    pill.fill.fore_color.rgb = ACCENT
    pill.line.fill.background()
    tx = slide.shapes.add_textbox(card_left + Pt(18), card_top + Pt(20), Pt(140), Pt(22))
    tf = tx.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = "MSIL-PASCO"
    p.font.name = FONT_NAME
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Big title
    title_box = slide.shapes.add_textbox(card_left + Pt(18), card_top + Pt(60), card_w - Pt(36), Pt(80))
    tf2 = title_box.text_frame
    tf2.clear()
    p2 = tf2.paragraphs[0]
    p2.text = spec.title
    p2.font.name = FONT_NAME
    p2.font.size = Pt(40)
    p2.font.bold = True
    p2.font.color.rgb = INK

    # Subtitle
    sub_lines = list(spec.subtitle_lines or [])
    if sub_lines:
        sub = slide.shapes.add_textbox(card_left + Pt(18), card_top + Pt(150), card_w - Pt(36), Pt(70))
        tf3 = sub.text_frame
        tf3.clear()
        for i, line in enumerate(sub_lines):
            p = tf3.paragraphs[0] if i == 0 else tf3.add_paragraph()
            p.text = line
            p.font.name = FONT_NAME
            p.font.size = Pt(18)
            p.font.color.rgb = MUTED
            p.space_after = Pt(2)

    _add_footer(slide, prs, "Repo: MSIL-PASCO  •  Deck auto-updated from codebase")


def _render_two_column_slide(slide, prs: Presentation, title: str, left_title: str, left_bullets: Sequence[str], right_title: str, right_bullets: Sequence[str]) -> None:
    _apply_theme_background(slide, prs)
    _clear_all_textboxes(slide)
    _add_header(slide, prs, title)
    _add_footer(slide, prs, "MSIL-PASCO")

    content_top = HEADER_H + Pt(28)
    content_h = prs.slide_height - HEADER_H - FOOTER_H - Pt(56)
    col_w = (prs.slide_width - (MARGIN_X * 2) - GUTTER) // 2

    left_x = MARGIN_X
    right_x = MARGIN_X + col_w + GUTTER

    _add_card(slide, left_x, content_top, col_w, content_h)
    _add_card(slide, right_x, content_top, col_w, content_h)

    _add_card_title(slide, left_x + Pt(16), content_top + Pt(16), col_w - Pt(32), left_title)
    _add_bullets_in_box(slide, left_x + Pt(16), content_top + Pt(44), col_w - Pt(32), content_h - Pt(60), left_bullets, bullet_symbol="◼")

    _add_card_title(slide, right_x + Pt(16), content_top + Pt(16), col_w - Pt(32), right_title)
    _add_bullets_in_box(slide, right_x + Pt(16), content_top + Pt(44), col_w - Pt(32), content_h - Pt(60), right_bullets, bullet_symbol="◼")


def _render_single_card_slide(slide, prs: Presentation, title: str, bullets: Sequence[str]) -> None:
    _apply_theme_background(slide, prs)
    _clear_all_textboxes(slide)
    _add_header(slide, prs, title)
    _add_footer(slide, prs, "MSIL-PASCO")

    content_top = HEADER_H + Pt(28)
    content_h = prs.slide_height - HEADER_H - FOOTER_H - Pt(56)
    card_w = prs.slide_width - (MARGIN_X * 2)
    _add_card(slide, MARGIN_X, content_top, card_w, content_h)
    _add_bullets_in_box(slide, MARGIN_X + Pt(18), content_top + Pt(18), card_w - Pt(36), content_h - Pt(36), bullets, bullet_symbol="•")


def _set_shape_text(
    shape, lines: Sequence[str], *, font_size: int, bold_first: bool = False
) -> None:
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.name = FONT_NAME
        p.font.size = Pt(font_size)
        if bold_first and i == 0:
            p.font.bold = True


def _find_placeholder(slide, placeholder_types: Sequence[int]):
    for shape in slide.shapes:
        if not getattr(shape, "is_placeholder", False):
            continue
        try:
            if shape.placeholder_format.type in placeholder_types:
                return shape
        except Exception:
            continue
    return None


def _ensure_title(slide):
    title_ph = _find_placeholder(
        slide, [PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE]
    )
    if title_ph is not None:
        return title_ph
    # Fallback: add a manual title textbox near top.
    left = Pt(36)
    top = Pt(18)
    width = slide.part.slide_layout.slide_master.part.presentation.slide_width - Pt(72)
    height = Pt(60)
    return slide.shapes.add_textbox(left, top, width, height)


def _find_body_placeholder(slide):
    body_ph = _find_placeholder(
        slide,
        [
            PP_PLACEHOLDER.BODY,
            PP_PLACEHOLDER.OBJECT,
            PP_PLACEHOLDER.SUBTITLE,
        ],
    )
    if body_ph is not None and getattr(body_ph, "has_text_frame", False):
        return body_ph
    # Fallback heuristic (first non-title text frame).
    title_shape = _ensure_title(slide)
    for shape in slide.shapes:
        if not getattr(shape, "has_text_frame", False):
            continue
        if shape is title_shape:
            continue
        return shape
    return None


def _set_title(slide, title: str) -> None:
    title_shape = _ensure_title(slide)
    tf = title_shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.level = 0
    p.font.name = FONT_NAME
    p.font.size = Pt(TITLE_SIZE)
    p.font.bold = True
    p.font.color.rgb = RGBColor(20, 30, 45)


def _set_bullets(slide, bullets: Sequence[str]) -> None:
    body = _find_body_placeholder(slide)
    if body is None:
        # Add a textbox if no body exists.
        left = Pt(60)
        top = Pt(120)
        width = slide.part.slide_layout.slide_master.part.presentation.slide_width - Pt(120)
        height = slide.part.slide_layout.slide_master.part.presentation.slide_height - Pt(160)
        body = slide.shapes.add_textbox(left, top, width, height)

    tf = body.text_frame
    tf.clear()
    tf.word_wrap = True

    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = b
        p.level = 0
        p.font.name = FONT_NAME
        p.font.size = Pt(BODY_SIZE)
        p.font.color.rgb = RGBColor(35, 45, 60)


def update_deck(input_pptx: Path, output_pptx: Path) -> None:
    prs = Presentation(str(input_pptx))

    specs: List[SlideSpec] = [
        SlideSpec(
            title="MSIL RVT – Mid Project Review",
            bullets=[],
            subtitle_lines=[
                "Real-time workshop visibility platform",
                "Eternal Robotics",
            ],
        ),
        SlideSpec(
            title="Project Overview",
            bullets=[
                "Provide real-time visibility into vehicle movement and workshop operations",
                "Automatically capture key events from IPC systems and applications",
                "Lay the foundation for operational metrics and analytics across stations",
            ],
        ),
        SlideSpec(
            title="What We Have Built So Far",
            bullets=[
                "Complete web UI with dashboards, vehicle tracking and monitoring views",
                "Backend services to capture events from IPC systems and applications",
                "DeepStream pipelines integrated for vehicle and bay-level tracking",
                "Technician tracking linked to vehicle and bay activities",
            ],
        ),
        SlideSpec(
            title="Solution Architecture (High Level)",
            bullets=[
                "Edge: cameras and DeepStream pipelines generate structured workshop events",
                "Backend: event APIs store visits, timelines and current workshop state",
                "Frontend: dashboards present station overview, vehicle flow and alerts",
                "Cloud: optional services for long-term storage and advanced analytics",
            ],
        ),
        SlideSpec(
            title="Current Capabilities (Without Cloud Dependencies)",
            bullets=[
                "Event capture and visit tracking run fully on local infrastructure",
                "Basic station overview metrics and live vehicle movement can be shown",
                "Technician movements can be associated with vehicles and bays",
                "Initial trends and operational patterns can be observed locally",
            ],
        ),
        SlideSpec(
            title="Next Steps",
            bullets=[
                "Reconnect cloud services for richer analytics and long-term trends",
                "Harden deployment, monitoring and operational readiness",
                "Expand reporting and KPIs for service advisors, reception and workshop leads",
                "Plan rollout and scale to additional stations and use cases",
            ],
        ),
    ]

    for idx, spec in enumerate(specs):
        if idx >= len(prs.slides):
            break
        slide = prs.slides[idx]

        # Fully re-render each slide for a cleaner, consistent UI
        if idx == 0:
            _render_title_slide(slide, prs, spec)
            continue

        # Purposefully vary layout to reduce "wall of text"
        if idx == 3:  # Architecture overview → 2 columns
            _render_two_column_slide(
                slide,
                prs,
                spec.title,
                "Event Ingestion",
                [
                    "Cameras and DeepStream pipelines generate structured events",
                    "Applications contribute additional workshop events",
                    "Events are normalised and persisted as visits and timelines",
                ],
                "State + Visualization",
                [
                    "Backend services keep track of active vehicles and workshop state",
                    "UI dashboards show station overview, vehicle flow and alerts",
                    "Cloud services add long-term storage and analytics when enabled",
                ],
            )
            continue

        _render_single_card_slide(slide, prs, spec.title, spec.bullets)

    output_pptx.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_pptx))


def main() -> None:
    input_pptx = Path(r"C:\Users\bhara\Downloads\msil_rvt_mid_project_review.pptx")
    output_pptx = Path(r"C:\Users\bhara\Downloads\msil_rvt_mid_project_review_MIDREVIEW_GENERIC.pptx")
    update_deck(input_pptx, output_pptx)
    print(f"Saved: {output_pptx}")


if __name__ == "__main__":
    main()

